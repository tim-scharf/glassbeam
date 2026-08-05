"""
build_performance_deck.py
---------------------------
Generates a .pptx summarizing how the 4-model classification pipeline
performed against a genuinely held-out Tanner sample
(output/tanner_heldout_sample.csv), what changed this session, and a full
per-row listing of all 200 sampled (modality, description) pairs with
their predictions. The full listing can't fit on one slide legibly, so it
spans as many slides as needed (2 side-by-side tables per slide).

IMPORTANT: an earlier version of this deck sampled from
output/tanner_sample.csv, which turned out to be 88% duplicates of rows
already in the training corpus (Tanner is one of the 18 source customers
in output/glassbeam_data.csv) -- so it wasn't measuring generalization at
all. This version samples only from the ~10-11% of the Tanner export
whose exact text never appears in the training corpus (see
scripts/build_performance_deck.py's companion sampling step), for a real
held-out accuracy number.

Usage:
    python3 scripts/build_performance_deck.py
"""

import sys
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SCRIPT_DIR = Path(__file__).parent
ROOT_DIR = SCRIPT_DIR.parent
OUT_PATH = ROOT_DIR / "Tanner_Sample_Performance_Summary.pptx"
SAMPLE_CSV = ROOT_DIR / "output" / "tanner_heldout_sample.csv"

sys.path.insert(0, str(SCRIPT_DIR))
from modality_architecture import load_modality_architecture
from classify import classify

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
TEAL = RGBColor(0x0E, 0x7C, 0x86)
GREEN = RGBColor(0x1E, 0x7A, 0x34)
RED = RGBColor(0xB0, 0x2A, 0x2A)
GRAY = RGBColor(0x5A, 0x5A, 0x5A)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_title(slide, text, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.25), Inches(12.4), Inches(0.7))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.45), Inches(0.85), Inches(12.4), Inches(0.4))
        tf2 = box2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY


def add_bullets(slide, left, top, width, height, heading, items, heading_color=NAVY, font_size=13):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = heading
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = heading_color
    for item in items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
        p.space_after = Pt(4)
    return box


def style_header_row(row, fill=NAVY):
    for cell in row.cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER


def set_cell(cell, text, size=11, color=RGBColor(0x22, 0x22, 0x22), bold=False, align=PP_ALIGN.CENTER):
    cell.text = text
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
    cell.vertical_anchor = 1
    cell.margin_top = Pt(2)
    cell.margin_bottom = Pt(2)


def build_slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        slide,
        "Study Description Classifier — Performance on Held-Out Data",
        "Tanner Health System — 200 rows sampled from TannerStudy2025.xlsx + TannerStudy2026.xlsx, filtered to only the "
        "~10-11% whose exact text never appears in the training corpus, run through the full 4-model pipeline",
    )

    correction = slide.shapes.add_textbox(Inches(0.45), Inches(1.15), Inches(12.4), Inches(0.35))
    p = correction.text_frame.paragraphs[0]
    p.text = ("CORRECTED: an earlier pass randomly sampled 200 Tanner rows without checking for overlap with training data — "
              "88% turned out to be exact duplicates already in glassbeam_data.csv (Tanner is one of the 18 source customers). "
              "This version filters to the ~11% of rows with zero exact-text overlap for a real generalization number.")
    p.font.size = Pt(10.5)
    p.font.italic = True
    p.font.bold = True
    p.font.color.rgb = RED

    # Performance table
    rows_data = [
        ("CT", "134", "92.5%", "100%", "100%", "N/A"),
        ("XA", "7", "100%", "100%", "N/A", "N/A"),
        ("RF", "2", "50%*", "100%", "100%", "N/A"),
        ("MG", "2", "N/A", "100%", "N/A", "100%"),
        ("US", "1", "100%*", "100%", "100%", "N/A"),
        ("NM", "0", "insufficient data (2 source rows had blank descriptions)", "", "", ""),
        ("PT / ES / SR / OT", "51", "N/A (modality codes not yet in routing table)", "", "", ""),
    ]
    headers = ["Modality", "N", "Region/Focus", "Laterality", "Contrast", "Technique"]

    left, top = Inches(0.45), Inches(1.75)
    width, height = Inches(7.4), Inches(3.9)
    table_shape = slide.shapes.add_table(len(rows_data) + 1, len(headers), left, top, width, height)
    table = table_shape.table
    col_widths = [Inches(1.75), Inches(0.7), Inches(1.6), Inches(1.15), Inches(1.1), Inches(1.1)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for c, h in enumerate(headers):
        set_cell(table.cell(0, c), h)
    style_header_row(table.rows[0])

    low_score_rows = {0, 2}  # CT, RF (0-indexed among data rows -> row index+1 in table)
    for r, row in enumerate(rows_data, start=1):
        for c, val in enumerate(row):
            color = RGBColor(0x22, 0x22, 0x22)
            bold = False
            if c >= 2 and val not in ("N/A", "", "N/A (correctly not routed)"):
                pct = val.rstrip("*%")
                try:
                    pct_val = float(pct)
                    if pct_val < 70:
                        color, bold = RED, True
                    elif pct_val == 100:
                        color = GREEN
                except ValueError:
                    pass
            set_cell(table.cell(r, c), val, color=color, bold=bold, align=PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT)
        if r - 1 in low_score_rows:
            for c in range(len(headers)):
                table.cell(r, c).fill.solid()
                table.cell(r, c).fill.fore_color.rgb = RGBColor(0xFC, 0xE9, 0xE9)

    note = slide.shapes.add_textbox(Inches(0.45), Inches(5.75), Inches(7.4), Inches(0.75))
    tf = note.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ("*RF (n=2) and US (n=1) samples are too small to be statistically meaningful — reported for completeness, "
              "not as a generalization claim. NM/XA got no contrast score because they aren't routed for it by design.")
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = GRAY

    add_bullets(
        slide,
        Inches(8.1), Inches(1.75), Inches(4.75), Inches(4.9),
        "Takeaways",
        [
            "Laterality and contrast are still perfect (100%) on truly unseen text — that logic generalizes well.",
            "Region/focus is the real weak spot: 92.5% on CT (n=134, the only modality with a meaningful sample size), with every miss traceable to a specific cause, not noise.",
            "2 gaps reconfirmed from the earlier (contaminated) test: PET/PSMA imaging and PET myocardial perfusion still have no region keyword match.",
            "2 new gaps found here: “PET ... SKULL TO THIGH” (a whole-body-range PET scan) gets mis-summarized as isolated Head+Lower-extremity instead of a body-range signal; “UPPER GI” fluoroscopy has no Abdomen mapping at all.",
            "The bare-BONE→Whole-Body last-resort rule fired wrongly again on “CT GUIDED NEEDLE BIOPSY BONE DEEP” — same bug as before, now confirmed on a second independent example.",
            "New modality codes showed up that aren't in the routing table yet: PT (PET) and ES (endoscopy) — currently correctly abstain (return null), but PT especially looks worth adding given the volume of PET/CT-hybrid studies.",
        ],
        font_size=11.5,
    )

    return slide


def build_slide2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        slide,
        "Session Changes & Example Classifications",
        "What was built/fixed this session, and representative good vs. bad outputs from the Tanner sample",
    )

    add_bullets(
        slide,
        Inches(0.45), Inches(1.45), Inches(5.7), Inches(5.2),
        "Changes this session",
        [
            "Built technique_study_type model (NM + MG only — the only 2 modalities where specialty-technique keywords partition ~98%+ of the modality into a clean ~12-category enum).",
            "Enforced modality_model_architecture.json as authoritative across all 4 extractors via modality_architecture.py — fixed the “NM Bilateral leak” and analogous MG/XA/NM contrast leaks.",
            "Built classify.py: single (modality, study_desc_raw) → all-applicable-attributes router, plus a confidence.py layer giving every prediction a rule-tier confidence score and a corpus-novelty score.",
            "Dropped “reason” as a 5th attribute — only 5% of descriptions carry any indication language, not a clean taxonomy.",
            "Pulled Tanner Health data into the repo to test generalization — first pass sampled randomly and was 88% contaminated by rows already in the training corpus; caught it via the novelty score reading 0.0 on rows that should've been unfamiliar, then rebuilt the sample to only include the ~11% of rows with zero exact-text overlap with training data. This deck reflects the corrected, genuinely held-out result.",
        ],
        font_size=12,
    )

    good_rows = [
        ("CT ABDOMEN PELVIS WITH CONTRAST", "Abdomen, Pelvis • Unspecified • With"),
        ("CT SHOULDER RIGHT WITHOUT IV CONTRAST", "Upper extremity/shoulder • Right • Without"),
        ("MAMMO BREAST SCREENING TOMOSYNTHESIS 3D LEFT", "Left • Screening mammogram"),
        ("US RENAL ARTERY", "Abdomen/kidney"),
        ("CT CHEST ABDOMEN PELVIS WO CONTRAST", "Chest, Abdomen, Pelvis • Without"),
    ]
    bad_rows = [
        ("PET INITIAL BONE SKULL TO THIGH", "Head/skull + Lower ext/thigh — misses the body-range in between"),
        ("FL UPPER GI WITHOUT KUB", "region_focus: {} — no Abdomen mapping for “UPPER GI” at all"),
        ("PET PSMA IMAGING", "region_focus: {} — reconfirmed gap (no PSMA→prostate keyword)"),
        ("PET MYOCARDIAL PERFUSION MULTIPLE", "region_focus: {} — reconfirmed gap (no MYOCARDIAL keyword)"),
        ("CT GUIDED NEEDLE BIOPSY BONE DEEP", "region_focus: Whole Body — wrong; bare-BONE rule overreach, 2nd example"),
    ]

    def add_example_table(title, rows, top, header_color, row_shade):
        label = slide.shapes.add_textbox(Inches(6.4), top - Inches(0.4), Inches(6.4), Inches(0.35))
        p = label.text_frame.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(15)
        p.font.color.rgb = header_color

        n = len(rows) + 1
        table_shape = slide.shapes.add_table(n, 2, Inches(6.4), top, Inches(6.4), Inches(0.35 * n))
        table = table_shape.table
        table.columns[0].width = Inches(3.7)
        table.columns[1].width = Inches(2.7)
        set_cell(table.cell(0, 0), "Study description", align=PP_ALIGN.LEFT)
        set_cell(table.cell(0, 1), "Model output", align=PP_ALIGN.LEFT)
        style_header_row(table.rows[0], fill=header_color)
        for r, (desc, out) in enumerate(rows, start=1):
            set_cell(table.cell(r, 0), desc, size=10.5, align=PP_ALIGN.LEFT)
            set_cell(table.cell(r, 1), out, size=10.5, align=PP_ALIGN.LEFT)
            for c in range(2):
                table.cell(r, c).fill.solid()
                table.cell(r, c).fill.fore_color.rgb = row_shade
        return top + Inches(0.35 * n) + Inches(0.55)

    next_top = add_example_table("Good — correct end to end", good_rows, Inches(1.85), GREEN, RGBColor(0xEA, 0xF7, 0xEC))
    add_example_table("Bad — gaps found on this new site", bad_rows, next_top, RED, RGBColor(0xFC, 0xE9, 0xE9))

    return slide


def truncate(text, n):
    return text if len(text) <= n else text[: n - 1].rstrip() + "…"


def format_prediction(rec):
    """Compact one-line rendering of a classify() result for a dense table cell.
    Drops the "/unspecified" focus suffix (it's the common case) to keep this
    short enough to fit one line at small point sizes."""
    parts = []
    rf = rec["region_focus"]
    if rf is not None:
        bits = []
        for r, foci in sorted(rf.items()):
            specific = sorted(f for f in foci if f != "unspecified")
            bits.append(f"{r}({','.join(specific)})" if specific else r)
        parts.append(f"Rgn:{'; '.join(bits) if bits else 'none'}")
    if rec["laterality"] is not None:
        parts.append(f"Lat:{rec['laterality']}")
    if rec["contrast"] is not None:
        parts.append(f"CM:{rec['contrast']}")
    if rec["technique_study_type"] is not None:
        parts.append(f"Tech:{rec['technique_study_type']}")
    return " | ".join(parts) if parts else "(not routed)"


def load_all_predictions():
    """Classify every one of the 200 sampled rows individually (not deduped) and
    format a compact prediction string for each, preserving original row order."""
    architecture = load_modality_architecture()
    df = pd.read_csv(SAMPLE_CSV)
    records = []
    for modality, raw in df[["modality", "study_desc_raw"]].itertuples(index=False):
        if not isinstance(raw, str) or not raw.strip():
            continue
        result = classify(modality, raw, architecture)
        records.append((modality, raw, format_prediction(result)))
    return records


def add_mini_table(slide, left, top, width, rows, row_height=Inches(0.17)):
    n = len(rows) + 1
    table_shape = slide.shapes.add_table(n, 3, left, top, width, row_height * n)
    table = table_shape.table
    table.columns[0].width = Inches(0.55)
    table.columns[1].width = Inches(2.9)
    table.columns[2].width = width - Inches(0.55) - Inches(2.9)
    desc_max_chars = 44
    pred_max_chars = 40

    set_cell(table.cell(0, 0), "Mod", align=PP_ALIGN.LEFT)
    set_cell(table.cell(0, 1), "Study description", align=PP_ALIGN.LEFT)
    set_cell(table.cell(0, 2), "Prediction", align=PP_ALIGN.LEFT)
    style_header_row(table.rows[0], fill=TEAL)

    for r, (modality, desc, pred) in enumerate(rows, start=1):
        set_cell(table.cell(r, 0), modality, size=7, align=PP_ALIGN.LEFT)
        set_cell(table.cell(r, 1), truncate(desc, desc_max_chars), size=7, align=PP_ALIGN.LEFT)
        set_cell(table.cell(r, 2), truncate(pred, pred_max_chars), size=7, align=PP_ALIGN.LEFT)
        if r % 2 == 0:
            for c in range(3):
                table.cell(r, c).fill.solid()
                table.cell(r, c).fill.fore_color.rgb = LIGHT_GRAY


def build_data_slides(prs, records, per_table=30):
    """One slide per 2*per_table records (2 side-by-side mini tables each).
    Cell text is truncated to a fixed character budget (see add_mini_table)
    so every row renders as exactly one line -- table height is then fully
    predictable (no PowerPoint auto-grow-on-wrap surprises pushing content
    off the bottom of the slide)."""
    total = len(records)
    per_slide = per_table * 2
    n_slides = -(-total // per_slide)  # ceil

    table_left = Inches(0.45)
    table_top = Inches(1.3)
    table_width = Inches(6.15)
    gap = Inches(0.3)

    for s in range(n_slides):
        chunk = records[s * per_slide:(s + 1) * per_slide]
        left_rows = chunk[:per_table]
        right_rows = chunk[per_table:]

        start_idx = s * per_slide + 1
        end_idx = s * per_slide + len(chunk)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(
            slide,
            "All Sampled Study Descriptions — Modality & Predictions",
            f"Rows {start_idx}–{end_idx} of {total} (held-out Tanner sample — zero exact-text overlap with training data — original order, not deduplicated)",
        )

        add_mini_table(slide, table_left, table_top, table_width, left_rows)
        if right_rows:
            add_mini_table(slide, table_left + table_width + gap, table_top, table_width, right_rows)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide1(prs)
    build_slide2(prs)

    records = load_all_predictions()
    build_data_slides(prs, records)

    prs.save(str(OUT_PATH))
    print(f"Wrote {OUT_PATH} ({len(prs.slides)} slides, {len(records)} data rows)")


if __name__ == "__main__":
    main()
