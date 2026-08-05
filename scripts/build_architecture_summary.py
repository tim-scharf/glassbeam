"""
build_architecture_summary.py
--------------------------------
Generates a 1-slide .pptx summarizing (1) the current pipeline structure,
(2) the feedback loop used to grow/correct it, and (3) the intention to
introduce embeddings for region/focus specifically -- the one attribute
that hasn't generalized as cleanly to fresh data as the other three.

Usage:
    python3 scripts/build_architecture_summary.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT_DIR = Path(__file__).parent.parent
OUT_PATH = ROOT_DIR / "Architecture_Feedback_Loop_Roadmap.pptx"

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
TEAL = RGBColor(0x0E, 0x7C, 0x86)
PLUM = RGBColor(0x6B, 0x2E, 0x7A)
GRAY = RGBColor(0x5A, 0x5A, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x22, 0x22, 0x22)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

COLUMN_TOP = Inches(1.55)
COLUMN_HEIGHT = Inches(5.55)
COLUMN_WIDTH = Inches(4.05)
COLUMN_GAP = Inches(0.2)
COLUMN_LEFTS = [Inches(0.45), Inches(0.45) + COLUMN_WIDTH + COLUMN_GAP, Inches(0.45) + 2 * (COLUMN_WIDTH + COLUMN_GAP)]


def add_title(slide):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.25), Inches(12.4), Inches(0.65))
    p = box.text_frame.paragraphs[0]
    p.text = "Study Description Classifier — Structure, Feedback Loop & Embeddings Roadmap"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY

    box2 = slide.shapes.add_textbox(Inches(0.45), Inches(0.82), Inches(12.4), Inches(0.4))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Rule-based, no-embeddings pipeline today; embeddings are the planned next step specifically for region/focus."
    p2.font.size = Pt(13)
    p2.font.color.rgb = GRAY


def add_column(slide, left, heading, heading_color, items, bold_last=False):
    header_box = slide.shapes.add_textbox(left, COLUMN_TOP, COLUMN_WIDTH, Inches(0.4))
    p = header_box.text_frame.paragraphs[0]
    p.text = heading
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = heading_color

    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, COLUMN_TOP + Inches(0.42), COLUMN_WIDTH, Pt(2.2))
    rule.fill.solid()
    rule.fill.fore_color.rgb = heading_color
    rule.line.fill.background()

    body_box = slide.shapes.add_textbox(left, COLUMN_TOP + Inches(0.58), COLUMN_WIDTH, COLUMN_HEIGHT - Inches(0.58))
    tf = body_box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = f"• {item}"
        p.font.size = Pt(10.5)
        p.font.color.rgb = INK
        p.space_after = Pt(6)


def build_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide)

    add_column(
        slide, COLUMN_LEFTS[0], "1 · Current structure", NAVY,
        [
            "Two inputs: modality + study_desc_raw (the raw study description text).",
            "modality_model_architecture.json is the routing table: which of the 4 attributes apply per modality (e.g. NM skips laterality; MG only gets laterality + technique).",
            "4 independent rule-based (regex/keyword, no embeddings) models, each its own extractor script:",
            "   – Region/Focus — multi-label, 9 fixed regions, growing per-region foci.",
            "   – Laterality — Right / Left / Bilateral / Unspecified.",
            "   – Contrast timing — With / Without / With and without / Unspecified.",
            "   – Technique/study type — NM + MG only (~12 cats each) — the only 2 modalities with a clean specialty-technique partition.",
            "classify.py is the single entry point: reads the routing table once, calls only the applicable models, returns one combined record.",
        ],
    )

    add_column(
        slide, COLUMN_LEFTS[1], "2 · Feedback loop", TEAL,
        [
            "Every keyword rule is derived from actually reading the real corpus (6,139 unique descriptions) — never guessed or templated.",
            "Validated against held-out Tanner Health rows. First attempt sampled randomly and was 88% duplicates of training data (Tanner is already one of the 18 source customers) — caught via novelty scoring 0.0 on rows that should've been unfamiliar; fixed by filtering to the ~11% with zero exact-text overlap. The loop catching its own contaminated test is itself the loop working.",
            "Failures triage into 3 buckets: missing keyword → patch vocabulary; structural asymmetry → fix the rule (e.g. technique_extractor knew MYOCARDIAL/HEPATOBILIARY but region_focus_extractor didn't); out of scope → drop it (e.g. “reason” — only 5% of descriptions carry indication language).",
            "Each new site's vocabulary folds back into the shared keyword sets, so coverage compounds over time without hand-labeled training data.",
            "Same loop applies to the routing table itself: entries get corrected when empirical spec-rates disagree (e.g. contrast removed from MG/XA/NM; laterality removed from NM).",
        ],
    )

    add_column(
        slide, COLUMN_LEFTS[2], "3 · Embeddings — planned for region/focus", PLUM,
        [
            "Region/focus is the weakest generalizer of the 4 on genuinely held-out data (92.5% on CT, the only modality with enough held-out volume to be meaningful) because it requires hand-authoring every synonym for 9 regions × growing foci (PSMA, MYOCARDIAL, UPPER GI, whole-body-range PET scans, …).",
            "Laterality, contrast, and technique are already ~100% on fresh data — rule-based is sufficient there; no embeddings planned for those.",
            "Plan: keep the rule-based region/focus model as the deterministic, auditable baseline — a fallback, not a throwaway.",
            "Add an embedding layer on top that maps novel phrasing to the same fixed 9-region/foci taxonomy by semantic similarity, instead of requiring an exact keyword match for every new term.",
            "The manually-audited fresh-data samples being accumulated now (Tanner, and future sites) double as the validation/anchor set for that embedding layer once built.",
        ],
    )

    return slide


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    build_slide(prs)
    prs.save(str(OUT_PATH))
    print(f"Wrote {OUT_PATH}")


if __name__ == "__main__":
    main()
